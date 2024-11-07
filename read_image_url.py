import os
import io
import base64
import json
import logging
import tempfile
import shutil
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
import requests

# Configure logging
logging.basicConfig(level=logging.INFO)

def extract_images(file_path):
    # Use TemporaryDirectory as context manager
    with tempfile.TemporaryDirectory() as temp_dir:
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        result = {
            "metadata": {
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "file_type": ext,
                "instructions": "The 'images' field contains a list of images extracted from the file, each with its metadata and URL."
            },
            "images": []
        }

        def save_image(image, image_name):
            image_path = os.path.join(temp_dir, f"{image_name}.png")
            image.save(image_path, format='PNG')
            return image_path

        def process_image_data(image_data, image_name, image_format):
            try:
                image = Image.open(io.BytesIO(image_data))
                image_path = save_image(image, image_name)
                # Correctly format the URL for Windows
                image_url = image_path.replace("\\", "/")
                result["images"].append({
                    "image_name": image_name,
                    "image_size": len(image_data),
                    "image_type": image_format.lower(),
                    "image_width": image.width,
                    "image_height": image.height,
                    "url": f"file:///{image_url}"  # URL to access the image
                })
                logging.info(f"Processed image: {image_name}")
            except Exception as e:
                logging.error(f"Error processing image {image_name}: {str(e)}")

        try:
            if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                with Image.open(file_path) as img:
                    image_name = os.path.splitext(os.path.basename(file_path))[0]
                    image_path = save_image(img, image_name)
                    image_url = image_path.replace("\\", "/")
                    result["images"].append({
                        "image_name": image_name,
                        "image_size": os.path.getsize(file_path),
                        "image_type": ext,
                        "image_width": img.width,
                        "image_height": img.height,
                        "url": f"file:///{image_url}"  # URL to access the image
                    })
                    logging.info(f"Processed image file: {file_path}")

            elif ext == '.pdf':
                doc = fitz.open(file_path)
                for page_num, page in enumerate(doc):
                    for img_index, img in enumerate(page.get_images(full=True)):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        image_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_page{page_num + 1}_img{img_index + 1}"
                        process_image_data(image_bytes, image_name, image_ext)

            elif ext in ['.docx', '.doc']:
                doc = Document(file_path)
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        image_data = rel.target_part.blob
                        image_name = rel.target_ref.split('/')[-1].split('.')[0]
                        process_image_data(image_data, image_name, image_name.split('.')[-1])

            elif ext in ['.pptx', '.ppt']:
                prs = Presentation(file_path)
                for slide_num, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            image_data = shape.image.blob
                            image_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_slide{slide_num + 1}_img{shape.shape_id}"
                            process_image_data(image_data, image_name, shape.image.ext)

            elif ext in ['.xlsx', '.xls']:
                wb = openpyxl.load_workbook(file_path, read_only=True)
                for sheet in wb.worksheets:
                    for image in sheet._images:
                        image_data = image._data()
                        image_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_{sheet.title}_img{image.anchor}"
                        process_image_data(image_data, image_name, image.format)

            elif ext == '.html':
                with open(file_path, 'r', encoding='utf-8') as file:
                    soup = BeautifulSoup(file, 'html.parser')
                    for img_tag in soup.find_all('img'):
                        img_src = img_tag.get('src')
                        if img_src.startswith('data:image'):
                            header, encoded = img_src.split(',', 1)
                            image_data = base64.b64decode(encoded)
                            image_ext = header.split('/')[1].split(';')[0]
                            image_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_img{len(result['images']) + 1}"
                            process_image_data(image_data, image_name, image_ext)
                        else:
                            response = requests.get(img_src)
                            if response.status_code == 200:
                                image_data = response.content
                                image_ext = img_src.split('.')[-1]
                                image_name = f"{os.path.splitext(os.path.basename(file_path))[0]}_img{len(result['images']) + 1}"
                                process_image_data(image_data, image_name, image_ext)

        except Exception as e:
            logging.error(f"Error processing file {file_path}: {str(e)}")
            result["metadata"]["error"] = f"Error processing file: {str(e)}"

        return result



