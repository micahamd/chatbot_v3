import os
import io
import re
from PIL import Image
import pytesseract
import fitz
from docx import Document
from pptx import Presentation
import csv
import openpyxl
import hashlib

def extract_json_text(file_path):
    def apply_ocr(image):
        try:
            return pytesseract.image_to_string(image)
        except pytesseract.TesseractNotFoundError:
            return ""

    def hash_content(content):
        return hashlib.md5(content.encode('utf-8')).hexdigest()[:8]  # Shortened hash

    def add_content(content, check_similar=True):
        if len(content.strip()) < 5:
            return None
            
        content_hash = hash_content(content)
        
        # Check for similar existing content
        if check_similar:
            for existing_hash, existing_content in unique_contents.items():
                if content.strip() in existing_content or existing_content in content.strip():
                    return existing_hash
                    
        if content_hash not in unique_contents:
            unique_contents[content_hash] = content
        return content_hash

    # Simplified metadata structure
    result = {
        "f": os.path.basename(file_path),
        "t": os.path.splitext(file_path)[1].lower(),
        "content": {
            "text": {},
            "pages": []
        }
    }

    unique_contents = {}

    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            image = Image.open(file_path)
            ocr_text = apply_ocr(image)
            if ocr_text.strip():
                content_hash = add_content(ocr_text.strip())
                if content_hash:
                    result["content"]["pages"].append({"n": 1, "h": [content_hash]})

        elif ext == '.pdf':
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc, start=1):
                # Get blocks of text to preserve layout structure
                blocks = page.get_text("blocks")
                page_content = []
                
                if blocks:  # If text blocks exist
                    for block in blocks:
                        block_text = block[4].strip()  # block[4] contains the text content
                        if block_text:
                            content_hash = add_content(block_text)
                            if content_hash:
                                page_content.append(content_hash)
                else:  # Fallback to OCR if no text blocks
                    pix = page.get_pixmap()
                    img = Image.open(io.BytesIO(pix.tobytes()))
                    ocr_text = apply_ocr(img)
                    if ocr_text.strip():
                        content_hash = add_content(ocr_text.strip())
                        if content_hash:
                            page_content.append(content_hash)
                
                if page_content:
                    result["content"]["pages"].append({
                        "n": page_num,
                        "h": page_content,
                        "type": "text",
                        "metadata": {
                            "rotation": page.rotation,
                            "size": {"width": page.rect.width, "height": page.rect.height}
                        }
                    })
            
            doc.close()

        elif ext in ['.docx', '.doc']:
            doc = Document(file_path)
            current_page = []
            page_num = 1
            
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    if element.text.strip():
                        content_hash = add_content(element.text.strip())
                        if content_hash:
                            current_page.append(content_hash)
                elif element.tag.endswith('sectPr'):  # Section/Page break
                    if current_page:
                        result["content"]["pages"].append({
                            "n": page_num,
                            "h": current_page,
                            "type": "text"
                        })
                        current_page = []
                        page_num += 1
            
            # Add remaining content
            if current_page:
                result["content"]["pages"].append({
                    "n": page_num,
                    "h": current_page,
                    "type": "text"
                })

        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        content_hash = add_content(shape.text.strip())
                        if content_hash:
                            slide_content.append(content_hash)
                if slide_content:
                    result["content"]["pages"].append({"n": slide_num, "h": slide_content})

        elif ext in ['.xlsx', '.xls']:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                sheet_content = []
                for row in sheet.iter_rows():
                    row_text = ' '.join(str(cell.value).strip() for cell in row if cell.value)
                    if row_text:
                        content_hash = add_content(row_text)
                        if content_hash:
                            sheet_content.append(content_hash)
                if sheet_content:
                    result["content"]["pages"].append({"s": sheet.title, "h": sheet_content})

        elif ext == '.csv':
            with open(file_path, 'r', encoding='utf-8-sig') as csvfile:
                reader = csv.reader(csvfile)
                page_content = []
                for row in reader:
                    row_text = ' '.join(cell.strip() for cell in row if cell.strip())
                    if row_text:
                        content_hash = add_content(row_text)
                        if content_hash:
                            page_content.append(content_hash)
                if page_content:
                    result["content"]["pages"].append({"n": 1, "h": page_content})

        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read().strip()
                if text:
                    content_hash = add_content(text)
                    if content_hash:
                        result["content"]["pages"].append({"n": 1, "h": [content_hash]})

    except Exception:
        pass  # Removed error messages

    result["content"]["text"] = unique_contents
    return result